#!/usr/bin/env python3
"""
html2pptx.py - Convert Amplifier Stories HTML decks to PowerPoint presentations.

This tool parses HTML presentation decks created with the Amplifier Stories style
and generates equivalent PowerPoint presentations using python-pptx.

Usage:
    uv run --with python-pptx,beautifulsoup4,lxml python tools/html2pptx.py <input.html> [output.pptx]

If output path is not specified, uses the input filename with .pptx extension.

The tool extracts:
- Slide structure (each .slide div)
- Headlines and subheads
- Section labels
- Cards (title + text)
- Tenet boxes (accent-bordered info boxes)
- Feature lists
- Tables
- Highlight/callout boxes
- Big numbers and stats
"""

import argparse
import re
import sys
from pathlib import Path
from typing import Optional

from bs4 import BeautifulSoup, Tag
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Color palette (matching Amplifier Stories style)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MS_BLUE = RGBColor(0x00, 0x78, 0xD4)
MS_CYAN = RGBColor(0x50, 0xE6, 0xFF)
MS_GREEN = RGBColor(0x00, 0xCC, 0x6A)
MS_ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
MS_RED = RGBColor(0xFF, 0x45, 0x3A)
GRAY_70 = RGBColor(0xB3, 0xB3, 0xB3)
GRAY_50 = RGBColor(0x80, 0x80, 0x80)
DARK_GRAY = RGBColor(0x1A, 0x1A, 0x1A)
BORDER_GRAY = RGBColor(0x33, 0x33, 0x33)


def parse_color_from_class(classes: list[str]) -> Optional[RGBColor]:
    """Extract accent color from CSS classes."""
    color_map = {
        "green": MS_GREEN,
        "orange": MS_ORANGE,
        "red": MS_RED,
        "ms-green": MS_GREEN,
        "ms-orange": MS_ORANGE,
        "ms-red": MS_RED,
        "ms-blue": MS_BLUE,
        "ms-cyan": MS_CYAN,
        "warning": MS_ORANGE,
    }
    for cls in classes:
        if cls in color_map:
            return color_map[cls]
    return None


def get_text(element: Optional[Tag]) -> str:
    """Extract text content from an element, handling None."""
    if element is None:
        return ""
    return element.get_text(strip=True)


def set_slide_background(slide, color=BLACK):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 14,
    bold: bool = False,
    italic: bool = False,
    color: RGBColor = WHITE,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    wrap: bool = True,
):
    """Add a text box with specified styling."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_section_label(slide, text: str, top: float = 0.6):
    """Add a blue uppercase section label."""
    return add_text_box(
        slide,
        text.upper(),
        left=0.8,
        top=top,
        width=8.4,
        height=0.4,
        font_size=14,
        bold=True,
        color=MS_BLUE,
        wrap=False,
    )


def add_headline(
    slide, text: str, top: float = 1.1, size: int = 48, center: bool = False, color: RGBColor = WHITE
):
    """Add a large headline."""
    return add_text_box(
        slide,
        text,
        left=0.8,
        top=top,
        width=8.4,
        height=1.5,
        font_size=size,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
    )


def add_subhead(slide, text: str, top: float = 2.8, color: RGBColor = GRAY_70, center: bool = False):
    """Add a subtitle/subheading."""
    return add_text_box(
        slide,
        text,
        left=0.8,
        top=top,
        width=8.4,
        height=1.0,
        font_size=24,
        color=color,
        align=PP_ALIGN.CENTER if center else PP_ALIGN.LEFT,
    )


def add_card(
    slide,
    title: str,
    text: str,
    left: float,
    top: float,
    width: float = 2.6,
    height: float = 1.8,
    title_color: RGBColor = MS_BLUE,
):
    """Add a card with title and description."""
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_GRAY
    card.line.color.rgb = BORDER_GRAY
    card.line.width = Pt(1)

    # Card title
    add_text_box(
        slide,
        title,
        left=left + 0.15,
        top=top + 0.15,
        width=width - 0.3,
        height=0.4,
        font_size=16,
        bold=True,
        color=title_color,
    )

    # Card text
    add_text_box(
        slide,
        text,
        left=left + 0.15,
        top=top + 0.5,
        width=width - 0.3,
        height=height - 0.6,
        font_size=12,
        color=GRAY_70,
    )


def add_tenet(
    slide,
    title: str,
    text: str,
    left: float,
    top: float,
    width: float = 4.0,
    height: float = 0.9,
    accent_color: RGBColor = MS_GREEN,
):
    """Add a tenet box with left border accent."""
    # Background color based on accent
    if accent_color == MS_GREEN:
        bg_color = RGBColor(0x0D, 0x1A, 0x0D)
    elif accent_color == MS_ORANGE:
        bg_color = RGBColor(0x1A, 0x15, 0x0D)
    elif accent_color == MS_RED:
        bg_color = RGBColor(0x1A, 0x0D, 0x0D)
    else:
        bg_color = RGBColor(0x0D, 0x15, 0x1A)

    # Background
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()

    # Left accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.05), Inches(height)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    # Title
    add_text_box(
        slide,
        title,
        left=left + 0.15,
        top=top + 0.1,
        width=width - 0.3,
        height=0.3,
        font_size=14,
        bold=True,
        color=WHITE,
    )

    # Text
    add_text_box(
        slide,
        text,
        left=left + 0.15,
        top=top + 0.4,
        width=width - 0.3,
        height=height - 0.5,
        font_size=11,
        color=GRAY_70,
    )


def add_highlight_box(slide, text: str, top: float = 4.2, color: RGBColor = MS_BLUE):
    """Add a highlight/callout box."""
    if color == MS_GREEN:
        bg_color = RGBColor(0x00, 0x1A, 0x0D)
    elif color == MS_ORANGE:
        bg_color = RGBColor(0x33, 0x1A, 0x00)
    else:
        bg_color = RGBColor(0x00, 0x1A, 0x33)

    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(8.4), Inches(0.7))
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = color
    box.line.width = Pt(1)

    add_text_box(
        slide,
        text,
        left=1.0,
        top=top + 0.15,
        width=8.0,
        height=0.5,
        font_size=14,
        color=WHITE,
    )


class HTMLToPPTXConverter:
    """Converts Amplifier Stories HTML decks to PowerPoint."""

    def __init__(self, html_content: str):
        self.soup = BeautifulSoup(html_content, "lxml")
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        self.blank_layout = self.prs.slide_layouts[6]

    def extract_slides(self) -> list[Tag]:
        """Extract all slide divs from the HTML."""
        return self.soup.find_all("div", class_="slide")

    def is_centered(self, slide_div: Tag) -> bool:
        """Check if slide has center class."""
        classes = slide_div.get("class", [])
        return "center" in classes

    def process_slide(self, slide_div: Tag, slide_num: int):
        """Process a single slide div and add to presentation."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        set_slide_background(slide)

        is_centered = self.is_centered(slide_div)
        current_top = 0.6

        # Extract section label
        section_label = slide_div.find(class_="section-label")
        if section_label:
            if is_centered:
                current_top = 1.5
            add_section_label(slide, get_text(section_label), top=current_top)
            current_top += 0.5

        # Extract headline (h1 or .headline)
        headline = slide_div.find(["h1", "h2"], class_="headline") or slide_div.find("h1")
        if headline:
            text = get_text(headline).replace("<br>", "\n")
            # Check for gradient/big text styling
            has_gradient = "big-text" in headline.get("class", [])
            color = MS_CYAN if has_gradient else WHITE
            size = 56 if headline.name == "h1" or has_gradient else 40

            if is_centered:
                current_top = max(current_top, 2.0)

            add_headline(slide, text, top=current_top, size=size, center=is_centered, color=color)
            current_top += 1.2 if size > 45 else 0.9

        # Extract medium headline (h2.medium-headline)
        medium_headline = slide_div.find(class_="medium-headline")
        if medium_headline and medium_headline != headline:
            add_headline(slide, get_text(medium_headline), top=current_top, size=36, center=is_centered)
            current_top += 0.8

        # Extract subhead
        subhead = slide_div.find(class_="subhead") or slide_div.find("p", class_="subhead")
        if subhead:
            text = get_text(subhead)
            add_subhead(slide, text, top=current_top, center=is_centered)
            current_top += 0.8

        # Extract cards (.card elements in .thirds or .halves or .fourths)
        card_containers = slide_div.find_all(class_=["thirds", "halves", "fourths"])
        for container in card_containers:
            cards = container.find_all(class_="card")
            if cards:
                self._add_cards(slide, cards, current_top)
                current_top += 2.0

        # Extract standalone cards not in containers
        standalone_cards = [
            c for c in slide_div.find_all(class_="card") if not c.find_parent(class_=["thirds", "halves", "fourths"])
        ]
        if standalone_cards:
            self._add_cards(slide, standalone_cards, current_top)
            current_top += 2.0

        # Extract tenet boxes
        tenets = slide_div.find_all(class_="tenet")
        if tenets:
            self._add_tenets(slide, tenets, current_top)
            current_top += len(tenets) * 0.5 + 0.5

        # Extract versus comparison
        versus = slide_div.find(class_="versus")
        if versus:
            self._add_versus(slide, versus, current_top)
            current_top += 2.5

        # Extract tables
        tables = slide_div.find_all("table", class_="data-table")
        for table in tables:
            self._add_table(slide, table, current_top)
            current_top += 2.5

        # Extract feature lists
        feature_lists = slide_div.find_all(class_="feature-list")
        for fl in feature_lists:
            if not fl.find_parent(class_="versus"):  # Skip lists inside versus
                self._add_feature_list(slide, fl, current_top)
                current_top += 1.5

        # Extract highlight boxes
        highlight_boxes = slide_div.find_all(class_="highlight-box")
        for hb in highlight_boxes:
            classes = hb.get("class", [])
            color = parse_color_from_class(classes) or MS_BLUE
            text = get_text(hb)
            add_highlight_box(slide, text, top=min(current_top, 4.5), color=color)
            current_top += 0.8

        # Extract stats grid
        stat_grid = slide_div.find(class_="stat-grid")
        if stat_grid:
            self._add_stats(slide, stat_grid, current_top)

        # Extract quote
        quote = slide_div.find(class_="quote")
        if quote:
            self._add_quote(slide, quote, current_top)

        # Extract small text at bottom
        small_text = slide_div.find(class_="small-text")
        if small_text:
            add_text_box(
                slide,
                get_text(small_text),
                left=0.8,
                top=4.8,
                width=8.4,
                height=0.4,
                font_size=14,
                color=GRAY_50,
                align=PP_ALIGN.CENTER if is_centered else PP_ALIGN.LEFT,
            )

    def _add_cards(self, slide, cards: list[Tag], top: float):
        """Add a row of cards to the slide."""
        num_cards = len(cards)
        if num_cards == 0:
            return

        # Calculate layout
        total_width = 8.4
        gap = 0.2
        card_width = (total_width - gap * (num_cards - 1)) / num_cards
        card_width = min(card_width, 2.8)  # Max width

        start_left = 0.8
        if num_cards <= 3:
            card_width = 2.6
            gap = 0.3
            start_left = 0.8

        for i, card in enumerate(cards):
            title_el = card.find(class_="card-title")
            text_el = card.find(class_="card-text")
            number_el = card.find(class_="card-number")

            title = get_text(title_el) if title_el else ""
            text = get_text(text_el) if text_el else ""

            left = start_left + i * (card_width + gap)

            if number_el:
                # Big number card
                number = get_text(number_el)
                self._add_number_card(slide, number, title, text, left, top, card_width)
            else:
                add_card(slide, title, text, left, top, width=card_width, height=1.8)

    def _add_number_card(
        self, slide, number: str, title: str, text: str, left: float, top: float, width: float
    ):
        """Add a card with a big number."""
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(1.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_GRAY
        card.line.color.rgb = BORDER_GRAY
        card.line.width = Pt(1)

        # Big number
        add_text_box(
            slide,
            number,
            left=left + 0.1,
            top=top + 0.1,
            width=width - 0.2,
            height=0.7,
            font_size=48,
            bold=True,
            color=MS_CYAN,
            align=PP_ALIGN.CENTER,
        )

        # Title
        add_text_box(
            slide,
            title,
            left=left + 0.1,
            top=top + 0.8,
            width=width - 0.2,
            height=0.3,
            font_size=14,
            bold=True,
            color=MS_BLUE,
            align=PP_ALIGN.CENTER,
        )

        # Text
        add_text_box(
            slide,
            text,
            left=left + 0.1,
            top=top + 1.1,
            width=width - 0.2,
            height=0.5,
            font_size=10,
            color=GRAY_70,
            align=PP_ALIGN.CENTER,
        )

    def _add_tenets(self, slide, tenets: list[Tag], top: float):
        """Add tenet boxes to the slide."""
        num_tenets = len(tenets)

        # Determine layout: 2 columns if 4+ tenets
        if num_tenets >= 4:
            for i, tenet in enumerate(tenets):
                col = i % 2
                row = i // 2
                self._add_single_tenet(slide, tenet, 0.8 + col * 4.5, top + row * 1.0, width=4.2)
        else:
            for i, tenet in enumerate(tenets):
                self._add_single_tenet(slide, tenet, 0.8, top + i * 1.0, width=8.4)

    def _add_single_tenet(self, slide, tenet: Tag, left: float, top: float, width: float):
        """Add a single tenet box."""
        title_el = tenet.find(class_="tenet-title")
        text_el = tenet.find(class_="tenet-text")

        title = get_text(title_el) if title_el else ""
        text = get_text(text_el) if text_el else ""

        classes = tenet.get("class", [])
        accent_color = parse_color_from_class(classes) or MS_GREEN

        add_tenet(slide, title, text, left, top, width=width, accent_color=accent_color)

    def _add_versus(self, slide, versus: Tag, top: float):
        """Add a versus comparison layout."""
        sides = versus.find_all(class_="versus-side")
        if len(sides) < 2:
            return

        left_side = sides[0]
        right_side = sides[1]

        # Left title
        left_title = left_side.find(class_="versus-title")
        if left_title:
            classes = left_title.get("class", [])
            color = parse_color_from_class(classes) or MS_ORANGE
            add_text_box(
                slide, get_text(left_title), left=0.8, top=top, width=4.0, height=0.4, font_size=24, bold=True, color=color
            )

        # Left items
        left_list = left_side.find(class_="feature-list")
        if left_list:
            items = left_list.find_all("li")
            for i, item in enumerate(items):
                text = get_text(item)
                # Check for check/x marks
                if "✓" in text or "check" in str(item):
                    color = MS_GREEN
                elif "✗" in text or "x-mark" in str(item):
                    color = MS_RED
                else:
                    color = WHITE
                add_text_box(
                    slide, text, left=0.8, top=top + 0.5 + i * 0.35, width=4.0, height=0.35, font_size=14, color=color
                )

        # VS divider
        add_text_box(
            slide, "vs", left=4.5, top=top + 1.2, width=1.0, height=0.5, font_size=32, bold=True, color=GRAY_50, align=PP_ALIGN.CENTER
        )

        # Right title
        right_title = right_side.find(class_="versus-title")
        if right_title:
            classes = right_title.get("class", [])
            color = parse_color_from_class(classes) or MS_GREEN
            add_text_box(
                slide, get_text(right_title), left=5.5, top=top, width=4.0, height=0.4, font_size=24, bold=True, color=color
            )

        # Right items
        right_list = right_side.find(class_="feature-list")
        if right_list:
            items = right_list.find_all("li")
            for i, item in enumerate(items):
                text = get_text(item)
                if "✓" in text or "check" in str(item):
                    color = MS_GREEN
                elif "✗" in text or "x-mark" in str(item):
                    color = MS_RED
                else:
                    color = WHITE
                add_text_box(
                    slide, text, left=5.5, top=top + 0.5 + i * 0.35, width=4.0, height=0.35, font_size=14, color=color
                )

    def _add_table(self, slide, table: Tag, top: float):
        """Add a data table to the slide."""
        rows = table.find_all("tr")
        if not rows:
            return

        row_height = 0.32
        for row_idx, row in enumerate(rows):
            cells = row.find_all(["th", "td"])
            is_header = row.find("th") is not None

            # Calculate column widths based on content
            num_cols = len(cells)
            col_widths = [8.0 / num_cols] * num_cols
            if num_cols >= 3:
                col_widths = [2.5, 2.5, 3.0][:num_cols]

            left = 0.8
            for col_idx, cell in enumerate(cells):
                text = get_text(cell)
                width = col_widths[col_idx] if col_idx < len(col_widths) else 2.0

                if is_header:
                    color = MS_BLUE
                    font_size = 12
                    bold = True
                else:
                    # First column is usually label
                    if col_idx == 0:
                        color = WHITE
                        bold = True
                    else:
                        color = GRAY_70
                        bold = False
                    font_size = 11

                    # Check for status indicators
                    if "✓" in text:
                        color = MS_GREEN
                    elif "✗" in text:
                        color = MS_RED
                    elif "~" in text:
                        color = MS_ORANGE

                add_text_box(
                    slide,
                    text,
                    left=left,
                    top=top + row_idx * row_height,
                    width=width,
                    height=row_height,
                    font_size=font_size,
                    bold=bold,
                    color=color,
                )
                left += width

    def _add_feature_list(self, slide, feature_list: Tag, top: float):
        """Add a feature list to the slide."""
        items = feature_list.find_all("li")
        for i, item in enumerate(items):
            text = get_text(item)

            # Determine color from content
            if "✓" in text:
                color = MS_GREEN
            elif "✗" in text:
                color = MS_RED
            else:
                color = WHITE

            add_text_box(
                slide, text, left=0.8, top=top + i * 0.4, width=8.4, height=0.4, font_size=16, color=color
            )

    def _add_stats(self, slide, stat_grid: Tag, top: float):
        """Add a stats grid to the slide."""
        stats = stat_grid.find_all(class_="stat")
        num_stats = len(stats)
        if num_stats == 0:
            return

        width_per_stat = 8.4 / num_stats
        start_left = 0.8 + (8.4 - width_per_stat * num_stats) / 2

        for i, stat in enumerate(stats):
            number_el = stat.find(class_="stat-number")
            label_el = stat.find(class_="stat-label")

            number = get_text(number_el) if number_el else ""
            label = get_text(label_el) if label_el else ""

            left = start_left + i * width_per_stat

            # Number
            add_text_box(
                slide,
                number,
                left=left,
                top=top,
                width=width_per_stat,
                height=0.6,
                font_size=40,
                bold=True,
                color=MS_CYAN,
                align=PP_ALIGN.CENTER,
            )

            # Label
            add_text_box(
                slide,
                label,
                left=left,
                top=top + 0.6,
                width=width_per_stat,
                height=0.4,
                font_size=12,
                color=GRAY_70,
                align=PP_ALIGN.CENTER,
            )

    def _add_quote(self, slide, quote: Tag, top: float):
        """Add a quote to the slide."""
        text = get_text(quote)
        add_text_box(
            slide,
            f'"{text}"',
            left=0.8,
            top=top,
            width=8.4,
            height=1.2,
            font_size=24,
            italic=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

        # Look for attribution
        attribution = quote.find_next_sibling(class_="quote-attribution")
        if attribution:
            add_text_box(
                slide,
                get_text(attribution),
                left=0.8,
                top=top + 1.2,
                width=8.4,
                height=0.3,
                font_size=14,
                color=GRAY_50,
                align=PP_ALIGN.CENTER,
            )

    def convert(self) -> Presentation:
        """Convert the HTML to a PowerPoint presentation."""
        slides = self.extract_slides()

        for i, slide_div in enumerate(slides):
            self.process_slide(slide_div, i + 1)

        return self.prs

    def save(self, output_path: str):
        """Save the presentation to a file."""
        self.prs.save(output_path)


def main():
    parser = argparse.ArgumentParser(
        description="Convert Amplifier Stories HTML decks to PowerPoint presentations."
    )
    parser.add_argument("input", help="Input HTML file path")
    parser.add_argument("output", nargs="?", help="Output PPTX file path (default: same name as input)")

    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {input_path}", file=sys.stderr)
        sys.exit(1)

    if args.output:
        output_path = Path(args.output)
    else:
        output_path = input_path.with_suffix(".pptx")

    print(f"Converting: {input_path}")
    print(f"Output: {output_path}")

    # Read HTML content
    html_content = input_path.read_text(encoding="utf-8")

    # Convert
    converter = HTMLToPPTXConverter(html_content)
    converter.convert()
    converter.save(str(output_path))

    print(f"Done! Created {output_path}")


if __name__ == "__main__":
    main()
